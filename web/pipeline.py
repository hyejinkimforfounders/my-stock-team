#!/usr/bin/env python3
"""웹앱 데이터 파이프라인 — 종목명/코드 → (재무·시세) 실데이터 → PPTX.

에이전트 팀의 데이터 소스를 그대로 사용합니다.
  - 재무: DART OpenAPI (fnlttSinglAcnt.json, 연결재무제표)
  - 시세: FinanceDataReader

LLM 없이 규칙 기반 코멘트로 리포트를 채웁니다(빠르고 키만 있으면 동작).
가드레일: 매수/매도 단정·목표가 미제시, 출처·기준일 병기, 학습용 명시.
"""
import os, ssl, certifi, functools, io, zipfile, urllib.request, urllib.parse, json
import datetime as dt
import xml.etree.ElementTree as ET

# ── macOS python.org SSL 우회 (CLAUDE.md 규칙) ──
ssl._create_default_https_context = functools.partial(
    ssl.create_default_context, cafile=certifi.where())
_SSL_CTX = ssl.create_default_context(cafile=certifi.where())

import FinanceDataReader as fdr

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))   # 프로젝트 루트
SKILL_DIR = os.path.join(BASE, ".claude", "skills", "report-pptx")

# build_pptx 디자인 헬퍼 재사용
# 배포(Vercel)용으로 web/build_pptx.py 동봉본을 우선 임포트, 없으면 스킬 원본 사용
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))   # web/
try:
    import build_pptx as bp
except ImportError:
    if SKILL_DIR not in sys.path:
        sys.path.insert(0, SKILL_DIR)
    import build_pptx as bp


# ──────────────────────────── 환경/키 ────────────────────────────
def _load_env():
    path = os.path.join(BASE, ".env.local")
    if not os.path.exists(path):
        return
    with open(path, encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            k, v = line.split("=", 1)
            os.environ.setdefault(k.strip(), v.strip())


_load_env()
DART_KEY = os.environ.get("DART_API_KEY") or os.environ.get("DART_KEY")


def _get_json(url):
    with urllib.request.urlopen(url, context=_SSL_CTX, timeout=30) as r:
        return json.loads(r.read().decode("utf-8"))


# ──────────────────────────── 종목 해석 (DART corpCode.xml) ────────────────────────────
_CORP_CACHE = None


def _corp_records():
    """상장사 매핑 → [(corp_code, corp_name, stock_code), ...].

    동봉된 web/corp_map.json을 우선 사용(서버리스 콜드스타트 최적화).
    없으면 DART corpCode.xml을 내려받아 파싱(로컬 갱신용).
    """
    global _CORP_CACHE
    if _CORP_CACHE is not None:
        return _CORP_CACHE

    cached = os.path.join(os.path.dirname(os.path.abspath(__file__)), "corp_map.json")
    if os.path.exists(cached):
        with open(cached, encoding="utf-8") as f:
            _CORP_CACHE = [tuple(r) for r in json.load(f)]
        return _CORP_CACHE

    if not DART_KEY:
        raise RuntimeError("DART_API_KEY 미설정 (.env.local 확인)")
    url = f"https://opendart.fss.or.kr/api/corpCode.xml?crtfc_key={DART_KEY}"
    with urllib.request.urlopen(url, context=_SSL_CTX, timeout=60) as r:
        blob = r.read()
    zf = zipfile.ZipFile(io.BytesIO(blob))
    xml = zf.read(zf.namelist()[0])
    root = ET.fromstring(xml)
    recs = []
    for el in root.iter("list"):
        recs.append((
            (el.findtext("corp_code") or "").strip(),
            (el.findtext("corp_name") or "").strip(),
            (el.findtext("stock_code") or "").strip(),
        ))
    _CORP_CACHE = recs
    return recs


def resolve(query):
    """종목명 또는 6자리 코드 → {name, ticker, corp_code}. 모호하면 ValueError."""
    query = query.strip()
    recs = _corp_records()
    listed = [r for r in recs if r[2] and r[2] != " " and len(r[2]) == 6]

    if query.isdigit() and len(query) == 6:        # 6자리 종목코드
        for corp_code, name, code in listed:
            if code == query:
                return {"name": name, "ticker": code, "corp_code": corp_code}
        raise ValueError(f"종목코드 {query}에 해당하는 상장사를 DART에서 찾지 못했습니다.")

    # 이름 매칭: 완전일치 우선, 없으면 부분일치
    exact = [r for r in listed if r[1] == query]
    cand = exact or [r for r in listed if query in r[1]]
    if not cand:
        raise ValueError(f"'{query}'에 해당하는 상장사를 찾지 못했습니다. (정확한 종목명/6자리 코드로 입력)")
    if len(cand) > 1 and not exact:
        names = ", ".join(sorted({r[1] for r in cand})[:6])
        raise ValueError(f"'{query}'와(과) 일치하는 종목이 여럿입니다: {names} … 정확히 입력해 주세요.")
    corp_code, name, code = cand[0]
    return {"name": name, "ticker": code, "corp_code": corp_code}


# ──────────────────────────── 시세 (FinanceDataReader) ────────────────────────────
def _won(x):
    return f"{x:,.0f}원"


def fetch_price(ticker, asof):
    """최근 약 380일 일봉 → 가격 요약. asof: datetime.date 기준일."""
    start = asof - dt.timedelta(days=400)
    df = fdr.DataReader(ticker, start.isoformat(), asof.isoformat())
    if df is None or df.empty:
        raise ValueError(f"FinanceDataReader에서 {ticker} 시세를 가져오지 못했습니다.")
    df = df.dropna(subset=["Close"])
    close = df["Close"]
    last = float(close.iloc[-1])
    last_date = close.index[-1].date()
    prev = float(close.iloc[-2]) if len(close) > 1 else last
    chg_pct = (last / prev - 1) * 100 if prev else 0.0

    def ma(n):
        return float(close.tail(n).mean()) if len(close) >= 1 else last

    hi52 = float(close.max()); lo52 = float(close.min())
    hi_date = close.idxmax().date(); lo_date = close.idxmin().date()
    vol = df["Volume"]
    last_vol = float(vol.iloc[-1])
    vol60 = float(vol.tail(60).mean()) if len(vol) else last_vol
    vol_ratio = (last_vol / vol60) if vol60 else 0.0

    return {
        "last": last, "last_date": last_date, "chg_pct": chg_pct,
        "ma5": ma(5), "ma20": ma(20), "ma60": ma(60),
        "hi52": hi52, "hi_date": hi_date, "lo52": lo52, "lo_date": lo_date,
        "last_vol": last_vol, "vol60": vol60, "vol_ratio": vol_ratio,
        "n_days": len(close), "close_series": close,
    }


def price_table(p):
    header = ["항목", "값"]
    rows = [
        ["종가", _won(p["last"]) + f"  ({p['chg_pct']:+.2f}%)"],
        ["5일 이평", _won(p["ma5"])],
        ["20일 이평", _won(p["ma20"])],
        ["60일 이평", _won(p["ma60"])],
        ["52주 최고", _won(p["hi52"]) + f"  ({p['hi_date']})"],
        ["52주 최저", _won(p["lo52"]) + f"  ({p['lo_date']})"],
        ["거래량(60일 평균 대비)", f"{p['vol_ratio']:.2f}x"],
    ]
    return header, rows


def price_comments(p):
    out = [f"종가 {_won(p['last'])} ({p['chg_pct']:+.2f}%) · "
           f"이평 5/20/60일 {p['ma5']:,.0f}/{p['ma20']:,.0f}/{p['ma60']:,.0f}원"]
    pos = []
    if p["last"] >= p["ma20"]:
        pos.append("20일선 위")
    else:
        pos.append("20일선 아래")
    if p["last"] >= p["ma60"]:
        pos.append("60일선 위")
    else:
        pos.append("60일선 아래")
    out.append(f"종가는 {' · '.join(pos)}에 위치합니다.")
    off_hi = (p["last"] / p["hi52"] - 1) * 100
    off_lo = (p["last"] / p["lo52"] - 1) * 100
    out.append(f"52주 최고 대비 {off_hi:+.1f}%, 52주 최저 대비 {off_lo:+.1f}% 수준입니다.")
    vr = p["vol_ratio"]
    tone = "활발" if vr >= 1.2 else ("한산" if vr <= 0.8 else "보통")
    out.append(f"최근 거래량은 60일 평균의 {vr:.2f}배로 {tone}한 편입니다.")
    return out


# ──────────────────────────── 재무 (DART) ────────────────────────────
def _to_jo(amount_str):
    """DART 금액 문자열(원) → 조원(float) 또는 None."""
    if not amount_str:
        return None
    s = amount_str.replace(",", "").strip()
    if s in ("", "-"):
        return None
    try:
        return int(s) / 1e12
    except ValueError:
        return None


def fetch_financials(corp_code, base_year):
    """fnlttSinglAcnt.json 1회 호출로 당기/전기/전전기 3개년 연결 재무 추출."""
    if not DART_KEY:
        return None
    # 최신 연간 보고서가 아직 없을 수 있어 base_year부터 한 해씩 내려가며 시도
    data = None
    used_year = None
    for y in (base_year, base_year - 1):
        url = ("https://opendart.fss.or.kr/api/fnlttSinglAcnt.json?"
               + urllib.parse.urlencode({
                   "crtfc_key": DART_KEY, "corp_code": corp_code,
                   "bsns_year": str(y), "reprt_code": "11011"}))
        try:
            j = _get_json(url)
        except Exception:
            continue
        if j.get("status") == "000" and j.get("list"):
            data, used_year = j["list"], y
            break
    if not data:
        return None

    # 연결(CFS) 손익계산서/재무상태표만 사용
    def pick(account_keys, sj):
        for it in data:
            if it.get("fs_div") != "CFS" or it.get("sj_div") != sj:
                continue
            nm = (it.get("account_nm") or "").replace(" ", "")
            if any(k in nm for k in account_keys):
                return (_to_jo(it.get("thstrm_amount")),
                        _to_jo(it.get("frmtrm_amount")),
                        _to_jo(it.get("bfefrmtrm_amount")))
        return (None, None, None)

    rev = pick(["매출액", "수익(매출액)", "영업수익"], "IS")
    op = pick(["영업이익"], "IS")
    net = pick(["당기순이익"], "IS")
    liab = pick(["부채총계"], "BS")
    equity = pick(["자본총계"], "BS")

    years = [used_year, used_year - 1, used_year - 2]   # 당기, 전기, 전전기

    def fmt(v):
        return f"{v:.1f}" if v is not None else "확인 불가"

    def pct(a, b):
        if a is None or b in (None, 0):
            return "확인 불가"
        return f"{a / b * 100:.1f}%"

    header = ["항목 (단위: 조원)"] + [f"{y}" for y in years]
    rows = [
        ["매출액"] + [fmt(v) for v in rev],
        ["영업이익"] + [fmt(v) for v in op],
        ["영업이익률"] + [pct(op[i], rev[i]) for i in range(3)],
        ["당기순이익"] + [fmt(v) for v in net],
        ["부채비율"] + [pct(liab[i], equity[i]) for i in range(3)],
    ]
    return {
        "header": header, "rows": rows, "year": used_year,
        "rev": rev, "op": op, "net": net, "liab": liab, "equity": equity,
    }


def fin_comments(f):
    out = []
    rev, op, net = f["rev"], f["op"], f["net"]
    if rev[0] is not None and rev[1] not in (None, 0):
        g = (rev[0] / rev[1] - 1) * 100
        out.append(f"{f['year']} 매출액은 {rev[0]:.1f}조원으로 전년 대비 {g:+.1f}% 변화했습니다.")
    if op[0] is not None and rev[0] not in (None, 0):
        out.append(f"영업이익은 {op[0]:.1f}조원, 영업이익률은 {op[0]/rev[0]*100:.1f}% 수준입니다.")
    if f["liab"][0] is not None and f["equity"][0] not in (None, 0):
        dr = f["liab"][0] / f["equity"][0] * 100
        tone = "낮은" if dr < 100 else "높은"
        out.append(f"부채비율은 {dr:.1f}%로 상대적으로 {tone} 편입니다.")
    if not out:
        out.append("DART에서 일부 계정을 확인하지 못했습니다(확인 불가 항목 참조).")
    return out


# ──────────────────────────── 리스크 / 종합의견 (규칙 기반) ────────────────────────────
def risk_bullets(p, f):
    out = []
    if p:
        dd = (p["last"] / p["hi52"] - 1) * 100
        if dd <= -20:
            out.append(f"52주 최고 대비 {dd:.1f}%로 큰 폭의 조정 구간에 있어 변동성 리스크가 있습니다.")
        if p["vol_ratio"] >= 1.5:
            out.append("최근 거래량이 평균을 크게 웃돌아 단기 변동성 확대 가능성이 있습니다.")
    if f:
        if f["liab"][0] is not None and f["equity"][0] not in (None, 0):
            dr = f["liab"][0] / f["equity"][0] * 100
            if dr >= 150:
                out.append(f"부채비율이 {dr:.1f}%로 재무 레버리지 부담이 있습니다.")
        if f["op"][0] is not None and f["op"][1] is not None and f["op"][0] < f["op"][1]:
            out.append("직전 연도 대비 영업이익이 감소해 수익성 둔화 리스크가 있습니다.")
    out.append("거시 환경(금리·환율·업황)에 따른 시장 리스크가 상존합니다.")
    return out


def opinion_bullets(p, f):
    """가치투자 관점 종합 — 매수/매도 단정·목표가 없이 근거까지만."""
    out = []
    if f and f["rev"][0] is not None:
        out.append("펀더멘털: 매출·이익 추이와 영업이익률, 부채비율로 본 재무 체력을 함께 점검했습니다.")
    if p:
        off_lo = (p["last"] / p["lo52"] - 1) * 100
        off_hi = (p["last"] / p["hi52"] - 1) * 100
        out.append(f"밸류에이션 참고: 현재가는 52주 저점 대비 {off_lo:+.1f}%, 고점 대비 {off_hi:+.1f}% 구간입니다.")
    out.append("안전마진 관점에서 내재가치 대비 가격 괴리와 재무 안정성을 함께 고려할 필요가 있습니다.")
    out.append("본 자료는 판단 근거 제시까지이며, 매수·매도 의견이나 목표가는 제시하지 않습니다.")
    return out


# ──────────────────────────── 차트 이미지 ────────────────────────────
def add_native_chart(slide, top, height, close_series):
    """python-pptx 네이티브 라인차트(종가+60일 이평). matplotlib 불필요(서버리스 친화)."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Pt

    # 약 60개 포인트로 다운샘플 (가독성·용량)
    n = len(close_series)
    step = max(1, n // 60)
    s = close_series.iloc[::step]
    ma60 = close_series.rolling(60).mean().iloc[::step]

    cats = [d.strftime("%y/%m") for d in s.index]
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("종가", [round(float(v)) for v in s.values])
    cd.add_series("60일 이평", [None if v != v else round(float(v)) for v in ma60.values])

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, bp.MARGIN_X, top, bp.USABLE_W, height, cd)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    for plot in chart.plots:
        plot.series[0].format.line.color.rgb = bp.NAVY
        plot.series[0].format.line.width = Pt(1.75)
        plot.series[1].format.line.color.rgb = bp.ORANGE
        plot.series[1].format.line.width = Pt(1.0)
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(8)
    return top + height


# ──────────────────────────── PPTX 조립 (build_pptx 헬퍼 재사용) ────────────────────────────
def _render_blocks_from(slide, blocks, start_top):
    """render_blocks를 지정한 y 좌표부터 렌더링(차트 아래 영역 채우기용)."""
    saved = bp.BODY_TOP
    try:
        bp.BODY_TOP = int(start_top)
        bp.render_blocks(slide, blocks)
    finally:
        bp.BODY_TOP = saved


def _section_slide(prs, page, title, blocks):
    s = bp.blank_slide(prs)
    bp.add_title(s, title)
    bp.render_blocks(s, blocks)
    bp.add_footer(s, page)
    return s


def build_report(query):
    """종목 입력 → PPTX 파일 경로 반환. 결과 메타 dict와 함께."""
    asof = dt.date.today()
    info = resolve(query)
    name, ticker, corp_code = info["name"], info["ticker"], info["corp_code"]

    # 데이터 수집 (일부 실패해도 리포트는 생성)
    price = None
    price_err = None
    try:
        price = fetch_price(ticker, asof)
    except Exception as e:
        price_err = str(e)

    fin = fetch_financials(corp_code, asof.year)

    # PPTX 조립
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = bp.SLIDE_W
    prs.slide_height = bp.SLIDE_H

    bp.build_cover(prs, name, asof.isoformat())

    # 2. 재무
    fin_blocks = []
    if fin:
        fin_blocks.append(("table", (fin["header"], fin["rows"])))
        fin_blocks.append(("source", f"(출처: DART, {fin['year']}/연간 · 연결 · 기준일 {fin['year']}-12-31)"))
        for c in fin_comments(fin):
            fin_blocks.append(("bullet", c))
    else:
        fin_blocks.append(("para", "DART 연결 재무 데이터를 확인하지 못했습니다 (확인 불가)."))
        fin_blocks.append(("source", "(출처: DART) — API 키 미설정 또는 해당 보고서 없음"))
    _section_slide(prs, 2, "재무", fin_blocks)

    # 3. 차트 (네이티브 라인차트 + 출처 + 코멘트)
    from pptx.util import Inches
    s3 = bp.blank_slide(prs)
    bp.add_title(s3, "차트")
    if price:
        y = add_native_chart(s3, Inches(1.45), Inches(3.1), price["close_series"])
        blocks = [("source", f"(출처: FinanceDataReader, 기준일: {price['last_date']})")]
        blocks += [("bullet", c) for c in price_comments(price)]
        # 차트 아래부터 렌더링되도록 임시 시작점 조정
        _render_blocks_from(s3, blocks, y + Inches(0.15))
    else:
        bp.render_blocks(s3, [("para", f"시세 데이터를 확인하지 못했습니다: {price_err or '확인 불가'}")])
    bp.add_footer(s3, 3)

    # 4. 리스크
    risk_blocks = [("bullet", b) for b in risk_bullets(price, fin)]
    _section_slide(prs, 4, "리스크", risk_blocks)

    # 5. 종합의견
    op_blocks = [("bullet", b) for b in opinion_bullets(price, fin)]
    op_blocks.append(("para", "본 분석은 교육·학습 목적이며 투자 권유가 아닙니다."))
    _section_slide(prs, 5, "종합의견", op_blocks)

    # 메모리로 저장(서버리스 호환) + 로컬은 best-effort 디스크 저장
    fname = f"{name}_{asof.isoformat()}.pptx"
    buf = io.BytesIO()
    prs.save(buf)
    data = buf.getvalue()

    path = None
    try:
        out_dir = os.path.join(BASE, "reports", "pptx")
        os.makedirs(out_dir, exist_ok=True)
        path = os.path.join(out_dir, fname)
        with open(path, "wb") as fh:
            fh.write(data)
    except Exception:
        path = None       # 읽기전용 FS(Vercel 등)면 디스크 저장 생략

    return {
        "path": path, "filename": fname, "data": data,
        "name": name, "ticker": ticker, "corp_code": corp_code,
        "asof": asof.isoformat(),
        "has_price": price is not None, "has_fin": fin is not None,
    }


if __name__ == "__main__":
    import sys
    q = sys.argv[1] if len(sys.argv) > 1 else "삼성전자"
    m = build_report(q)
    print({k: (f"<{len(v)} bytes>" if k == "data" else v) for k, v in m.items()})
