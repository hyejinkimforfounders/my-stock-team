#!/usr/bin/env python3
"""report-pptx — 종목 리서치 Markdown을 투자증권 리서치 톤 PPTX로 변환.

입력 : reports/{종목명}.md
출력 : reports/pptx/{종목명}.pptx   (pptx 폴더 없으면 생성)

슬라이드 순서(고정):
  표지 → 종목 개요 → 재무 요약 → 가격/추세 → 뉴스·심리 → 리스크 → 한 줄 종합

사용:
  python3 build_pptx.py 삼성전자
  python3 build_pptx.py reports/삼성전자.md --date 2026-06-16
"""
import sys, os, re, argparse, datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml

# ──────────────────────────── 디자인 토큰 ────────────────────────────
NAVY      = RGBColor(0x1F, 0x2A, 0x44)   # 딥네이비 — 제목/표 헤더
ORANGE    = RGBColor(0xF3, 0x73, 0x21)   # 한화 오렌지 — 상단바/구분선/키수치
BODY_GRAY = RGBColor(0x40, 0x40, 0x40)   # 본문 그레이
SUB_GRAY  = RGBColor(0x8A, 0x8A, 0x8A)   # 출처/푸터 그레이
LINE_GRAY = "D9D9D9"                       # 표 괘선(얇은 그레이)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
FONT      = "맑은 고딕"                     # 한글 폰트 단일 고정

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)     # 16:9
MARGIN_X         = Inches(0.6)
USABLE_W         = SLIDE_W - 2 * MARGIN_X
BODY_TOP         = Inches(1.45)
BODY_BOTTOM      = Inches(6.95)                     # 푸터 위 한계선
ROW_H            = Inches(0.34)                     # 표 행 기준 높이

DISCLAIMER = "학습용 리서치 · 투자 권유 아님 (특정일 매매 타이밍·목표가 미제시)"

# 가드레일: 매수/매도·목표가 단정 표현(있으면 생성 중단)
BANNED = [
    "매수 추천", "매도 추천", "매수의견", "매도의견", "매수 의견", "매도 의견",
    "강력매수", "강력 매수", "적극매수", "목표주가", "목표가", "적정주가",
    "비중확대", "비중축소", "사라고", "팔라고", "전량매도", "분할매수 추천",
]


# ──────────────────────────── 폰트 헬퍼 ────────────────────────────
def style_run(run, size, color, bold=False):
    """run 폰트 = 맑은 고딕(라틴/한글 동일), 크기/색/굵기 지정."""
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def add_runs(paragraph, text, size, color):
    """인라인 **굵게** → 오렌지 볼드(키 수치 강조), 나머지는 지정 색."""
    parts = re.split(r"\*\*(.+?)\*\*", text)
    for i, seg in enumerate(parts):
        if seg == "":
            continue
        run = paragraph.add_run()
        run.text = seg
        if i % 2 == 1:                       # ** ** 안쪽
            style_run(run, size, ORANGE, bold=True)
        else:
            style_run(run, size, color)


# ──────────────────────────── 슬라이드 골격 ────────────────────────────
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text, size, color,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    style_run(run, size, color, bold)
    return tb


def add_footer(slide, page_no):
    add_text(slide, MARGIN_X, Inches(7.06), Inches(9), Inches(0.3),
             DISCLAIMER, 9, SUB_GRAY)
    add_text(slide, SLIDE_W - Inches(1.3), Inches(7.06), Inches(0.7), Inches(0.3),
             str(page_no), 9, SUB_GRAY, align=PP_ALIGN.RIGHT)


def add_title(slide, title):
    """본문 슬라이드 제목(딥네이비) + 하단 오렌지 구분선."""
    add_text(slide, MARGIN_X, Inches(0.45), USABLE_W, Inches(0.7),
             title, 26, NAVY, bold=True)
    add_rect(slide, MARGIN_X, Inches(1.2), USABLE_W, Pt(2.5), ORANGE)


# ──────────────────────────── 표지 ────────────────────────────
def build_cover(prs, name, date_str):
    s = blank_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.28), ORANGE)        # 상단 오렌지 바
    add_text(s, MARGIN_X, Inches(2.7), USABLE_W, Inches(1.4),
             name, 46, NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_rect(s, SLIDE_W/2 - Inches(0.7), Inches(4.0), Inches(1.4), Pt(3), ORANGE)
    add_text(s, MARGIN_X, Inches(4.25), USABLE_W, Inches(0.5),
             f"작성일 {date_str}", 16, BODY_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, MARGIN_X, Inches(6.7), USABLE_W, Inches(0.4),
             DISCLAIMER, 11, SUB_GRAY, align=PP_ALIGN.CENTER)
    return s


# ──────────────────────────── 표 렌더링 ────────────────────────────
def _set_cell_borders(cell, hexcolor=LINE_GRAY, w_pt=0.75):
    tcPr = cell._tc.get_or_add_tcPr()
    w = int(w_pt * 12700)
    for tag in ("lnL", "lnR", "lnT", "lnB"):
        for el in tcPr.findall(qn("a:" + tag)):
            tcPr.remove(el)
    for tag in ("lnB", "lnT", "lnR", "lnL"):     # 역순 insert → 최종 L,R,T,B
        el = parse_xml(
            f'<a:{tag} {nsdecls("a")} w="{w}" cap="flat" cmpd="sng" algn="ctr">'
            f'<a:solidFill><a:srgbClr val="{hexcolor}"/></a:solidFill>'
            f'<a:prstDash val="solid"/></a:{tag}>')
        tcPr.insert(0, el)


def _fill_cell(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def _write_cell(cell, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.06)
    cell.margin_top = cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if bold:                                     # 헤더: 단색
        run = p.add_run(); run.text = text
        style_run(run, size, color, bold=True)
    else:                                         # 본문: **강조** 허용
        add_runs(p, text, size, color)


def render_table(slide, top, header, rows):
    """헤더행=네이비/흰글자, 본문행=흰배경/그레이 괘선. 높이 넘치면 행 절단."""
    ncol = len(header)
    fsize = 12 if ncol <= 4 else (11 if ncol <= 6 else 10)

    avail = BODY_BOTTOM - top
    max_rows = int(avail / ROW_H)                # 헤더 포함 최대 행수
    truncated = 0
    if 1 + len(rows) > max_rows:
        keep = max(1, max_rows - 1)
        truncated = len(rows) - keep
        rows = rows[:keep]

    nrow = 1 + len(rows)
    tbl_h = int(ROW_H) * nrow
    gtbl = slide.shapes.add_table(nrow, ncol, MARGIN_X, top, USABLE_W, tbl_h).table
    gtbl.first_row = False
    gtbl.horz_banding = False

    # 열 너비: 1열(라벨) 1.5배, 나머지 균등
    weights = [1.5] + [1.0] * (ncol - 1)
    tot = sum(weights)
    for i, wt in enumerate(weights):
        gtbl.columns[i].width = int(USABLE_W * wt / tot)

    for c in range(ncol):
        gtbl.rows[0].height = int(ROW_H)
        cell = gtbl.cell(0, c)
        _fill_cell(cell, NAVY)
        _set_cell_borders(cell)
        _write_cell(cell, header[c], fsize, WHITE, bold=True,
                    align=PP_ALIGN.CENTER)

    for r, row in enumerate(rows, start=1):
        gtbl.rows[r].height = int(ROW_H)
        for c in range(ncol):
            val = row[c] if c < len(row) else ""
            cell = gtbl.cell(r, c)
            _fill_cell(cell, WHITE)
            _set_cell_borders(cell)
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            _write_cell(cell, val, fsize, BODY_GRAY, align=align)

    bottom = top + tbl_h
    if truncated:
        add_text(slide, MARGIN_X, bottom + Inches(0.05), USABLE_W, Inches(0.3),
                 f"…외 {truncated}개 행 생략 (전체 수치는 원본 .md 참조)",
                 9, SUB_GRAY)
        bottom += Inches(0.35)
    return bottom


# ──────────────────────────── 본문 블록 렌더링 ────────────────────────────
def render_blocks(slide, blocks):
    """blocks: ('para'|'bullet'|'table'|'image'|'source', payload) 리스트."""
    y = BODY_TOP
    for kind, payload in blocks:
        if y >= BODY_BOTTOM:
            break
        if kind == "para":
            h = Inches(0.32) + Inches(0.22) * (len(payload) // 48)
            tb = slide.shapes.add_textbox(MARGIN_X, y, USABLE_W, h)
            tf = tb.text_frame; tf.word_wrap = True
            add_runs(tf.paragraphs[0], payload, 14, BODY_GRAY)
            y += h + Inches(0.06)
        elif kind == "bullet":
            h = Inches(0.30)
            tb = slide.shapes.add_textbox(MARGIN_X, y, USABLE_W, h)
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            lead = p.add_run(); lead.text = "•  "
            style_run(lead, 14, ORANGE, bold=True)
            add_runs(p, payload, 14, BODY_GRAY)
            y += h + Inches(0.02)
        elif kind == "source":
            tb = slide.shapes.add_textbox(MARGIN_X, y, USABLE_W, Inches(0.26))
            tf = tb.text_frame; tf.word_wrap = True
            run = tf.paragraphs[0].add_run(); run.text = payload
            style_run(run, 10, SUB_GRAY)
            run.font.italic = True
            y += Inches(0.30)
        elif kind == "table":
            header, rows = payload
            y = render_table(slide, y, header, rows) + Inches(0.12)
        elif kind == "image":
            y = render_image(slide, y, payload)
    return y


def render_image(slide, top, path):
    real = _resolve(path)
    maxw = USABLE_W
    maxh = BODY_BOTTOM - top
    if real and os.path.exists(real):
        pic = slide.shapes.add_picture(real, MARGIN_X, top)
        ratio = min(maxw / pic.width, maxh / pic.height, 1.0)
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
        pic.left = int(MARGIN_X + (USABLE_W - pic.width) / 2)
        return top + pic.height + Inches(0.1)
    # 차트 없음 → 안내 박스(수치 날조 금지)
    box = add_rect(slide, MARGIN_X, top, USABLE_W, Inches(2.2),
                   RGBColor(0xF2, 0xF2, 0xF2))
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    add_text(slide, MARGIN_X, top + Inches(0.95), USABLE_W, Inches(0.4),
             f"[차트 이미지 없음: {path}]", 12, SUB_GRAY, align=PP_ALIGN.CENTER)
    return top + Inches(2.3)


def _resolve(path):
    if os.path.isabs(path) and os.path.exists(path):
        return path
    for base in (PROJECT_ROOT, os.path.join(PROJECT_ROOT, "reports")):
        cand = os.path.join(base, path)
        if os.path.exists(cand):
            return cand
    return path


# ──────────────────────────── Markdown 파싱 ────────────────────────────
def parse_md(text):
    """H1=종목명, '작성일:'=날짜, ## 섹션별 블록 리스트 반환."""
    name, date_str, sections, cur = None, None, {}, None
    lines = text.splitlines()
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        if re.match(r"^#\s+", line):
            name = re.sub(r"^#\s+", "", line).strip()
            i += 1; continue
        m = re.match(r"^\s*작성일[:：]\s*(.+)$", line)
        if m and date_str is None:
            date_str = m.group(1).strip(); i += 1; continue
        if re.match(r"^##\s+", line):
            cur = re.sub(r"^##\s+", "", line).strip()
            sections[cur] = []
            i += 1; continue
        if cur is None:
            i += 1; continue
        blocks = sections[cur]
        if line.strip().startswith("|") and "|" in line.strip()[1:]:
            tbl_lines = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                tbl_lines.append(lines[i].strip()); i += 1
            header, rows = _parse_table(tbl_lines)
            if header:
                blocks.append(("table", (header, rows)))
            continue
        mimg = re.match(r"^\s*(?:!\[[^\]]*\]\(([^)]+)\)|chart[:：]\s*(\S+))", line)
        if mimg:
            blocks.append(("image", mimg.group(1) or mimg.group(2)))
            i += 1; continue
        if re.match(r"^\s*[-*]\s+", line):
            blocks.append(("bullet", re.sub(r"^\s*[-*]\s+", "", line).strip()))
            i += 1; continue
        s = line.strip()
        if s:
            if s.startswith("(출처") or s.startswith("출처"):
                blocks.append(("source", s))
            else:
                blocks.append(("para", s))
        i += 1
    return name, date_str, sections


def _parse_table(tbl_lines):
    rows = []
    for ln in tbl_lines:
        cells = [c.strip() for c in ln.strip().strip("|").split("|")]
        rows.append(cells)
    if len(rows) >= 2 and re.match(r"^[:\-\s|]+$", "|".join(rows[1])):
        return rows[0], rows[2:]
    return (rows[0], rows[1:]) if rows else (None, [])


# 섹션 헤더 → 고정 슬라이드 매핑(키워드)
SLIDE_SPEC = [
    ("종목 개요",  ["개요", "기업", "회사"]),
    ("재무 요약",  ["재무", "실적"]),
    ("가격 / 추세", ["가격", "추세", "주가", "차트", "기술"]),
    ("뉴스 · 심리", ["뉴스", "심리", "센티", "수급"]),
    ("리스크",     ["리스크", "위험"]),
    ("종합 의견",  ["종합", "결론", "의견"]),
]


def match_section(sections, keywords):
    for head, blocks in sections.items():
        if any(k in head for k in keywords):
            return blocks
    return None


# ──────────────────────────── 가드레일 검증 ────────────────────────────
def check_guardrails(text):
    errors, warnings = [], []
    low = text
    for term in BANNED:
        if term in low:
            errors.append(term)
    # 섹션별: 수치 있는데 출처 표기 없음 → 경고
    for blk in re.split(r"^##\s+", text, flags=re.M)[1:]:
        head = blk.splitlines()[0].strip()
        has_num = re.search(r"\d{2,}[%억조원]|\d+[.,]\d+|\d{4,}", blk)
        if has_num and "출처" not in blk:
            warnings.append(f"'{head}' 섹션에 수치가 있으나 (출처) 표기 없음")
    return errors, warnings


# ──────────────────────────── 메인 ────────────────────────────
def build(md_path, date_override=None):
    with open(md_path, encoding="utf-8") as f:
        raw = f.read()

    errors, warnings = check_guardrails(raw)
    if errors:
        print("⛔ 가드레일 위반 — 매수/매도·목표가 단정 표현이 있어 생성을 중단합니다:")
        for e in sorted(set(errors)):
            print(f"   · '{e}'")
        print("   → 판단 근거까지만 남기도록 .md를 수정한 뒤 다시 실행하세요.")
        sys.exit(2)
    for w in warnings:
        print(f"⚠️  {w}")

    name, date_str, sections = parse_md(raw)
    if not name:
        print("⛔ 종목명(H1 '# 종목명')을 찾지 못했습니다."); sys.exit(2)
    date_str = date_override or date_str or datetime.date.today().isoformat()

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    build_cover(prs, name, date_str)
    page = 2
    for title, keywords in SLIDE_SPEC:
        s = blank_slide(prs)
        add_title(s, title)
        blocks = match_section(sections, keywords)
        if blocks:
            render_blocks(s, blocks)
        else:
            add_text(s, MARGIN_X, BODY_TOP, USABLE_W, Inches(0.4),
                     "(해당 섹션 원본 내용 없음)", 13, SUB_GRAY)
        if title == "종합 의견":
            add_text(s, MARGIN_X, Inches(6.55), USABLE_W, Inches(0.4),
                     DISCLAIMER, 11, SUB_GRAY)
        add_footer(s, page)
        page += 1

    out_dir = os.path.join(PROJECT_ROOT, "reports", "pptx")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"{name}.pptx")
    prs.save(out_path)
    print(f"✅ 생성 완료: {os.path.relpath(out_path, PROJECT_ROOT)}  (슬라이드 {page-1}장)")
    return out_path


def main():
    ap = argparse.ArgumentParser(description="종목 리서치 .md → 리서치 톤 .pptx")
    ap.add_argument("target", help="종목명 또는 reports/경로.md")
    ap.add_argument("--date", help="작성일(YYYY-MM-DD), 미지정 시 .md의 작성일 또는 오늘")
    a = ap.parse_args()

    if a.target.endswith(".md"):
        md = a.target if os.path.isabs(a.target) else os.path.join(PROJECT_ROOT, a.target)
    else:
        md = os.path.join(PROJECT_ROOT, "reports", f"{a.target}.md")
    if not os.path.exists(md):
        print(f"⛔ 입력 파일 없음: {md}"); sys.exit(2)
    build(md, a.date)


PROJECT_ROOT = os.path.abspath(
    os.path.join(os.path.dirname(__file__), "..", "..", ".."))

if __name__ == "__main__":
    main()
